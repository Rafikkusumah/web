from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - profesional untuk 50+
DARK_RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x1F, 0x29, 0x37)
MEDIUM_GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
VERY_LIGHT_GRAY = RGBColor(0xFA, 0xFA, 0xFA)
RED_LIGHT = RGBColor(0xFE, 0xF2, 0xF2)

def add_background(slide, color):
    """Add solid background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add rectangle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add textbox with text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_accent_bar(slide, left, top, width, height, color=DARK_RED):
    """Add accent bar"""
    add_shape(slide, left, top, width, height, fill_color=color)

def add_circle(slide, left, top, size, fill_color):
    """Add circle shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ============================================
# SLIDE 1: COVER / TITLE SLIDE
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Top red bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), fill_color=DARK_RED)

# Main content area with subtle background
add_shape(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(5), fill_color=RED_LIGHT)

# Left accent bar
add_shape(slide, Inches(1), Inches(1.5), Inches(0.15), Inches(5), fill_color=DARK_RED)

# Company name - large and prominent
add_textbox(slide, Inches(2), Inches(2), Inches(9), Inches(1.2), 
            "PT CAHAYA DIMENSI BUMI", 44, True, DARK_RED, PP_ALIGN.LEFT)

# Divider line
add_shape(slide, Inches(2), Inches(3.3), Inches(3), Inches(0.05), fill_color=DARK_RED)

# Subtitle
add_textbox(slide, Inches(2), Inches(3.6), Inches(9), Inches(0.8), 
            "General Contractor & Automatic Door Solutions", 28, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Description
add_textbox(slide, Inches(2), Inches(4.5), Inches(9), Inches(1.2), 
            "Solusi Konstruksi Profesional & Sistem Pintu Otomatis dormakaba\nuntuk Kebutuhan Komersial dan Industri di Indonesia", 20, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Date and presenter info
add_textbox(slide, Inches(2), Inches(6), Inches(9), Inches(0.5), 
            f"Company Profile Presentation | {datetime.datetime.now().strftime('%B %Y')}", 16, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# ============================================
# SLIDE 2: COMPANY PROFILE
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), 
            "TENTANG PERUSAHAAN", 32, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1), Inches(0.08), WHITE)

# Left column - Main description
add_shape(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8), fill_color=VERY_LIGHT_GRAY)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(0.12), Inches(4.8), DARK_RED)

add_textbox(slide, Inches(1.3), Inches(2), Inches(6.5), Inches(0.6), 
            "Siapa Kami", 28, True, DARK_GRAY, PP_ALIGN.LEFT)

add_textbox(slide, Inches(1.3), Inches(2.7), Inches(6.5), Inches(3.5), 
            "PT Cahaya Dimensi Bumi adalah perusahaan general kontraktor yang berfokus pada solusi konstruksi berkualitas tinggi dan merupakan partner resmi dormakaba untuk solusi pintu otomatis di Indonesia.\n\n"
            "Dengan pengalaman dan keahlian yang mumpuni, kami berkomitmen memberikan pelayanan terbaik dalam setiap proyek konstruksi, mulai dari pembangunan baru, renovasi, hingga instalasi sistem pintu otomatis dormakaba.\n\n"
            "Kami telah dipercaya oleh lebih dari 100 perusahaan di Indonesia dan telah menyelesaikan berbagai proyek di sektor hospitality, industri, retail, dan infrastruktur.", 
            20, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Right column - Key facts
add_shape(slide, Inches(8.8), Inches(1.8), Inches(3.8), Inches(4.8), fill_color=DARK_RED)

add_textbox(slide, Inches(9.2), Inches(2.1), Inches(3), Inches(0.5), 
            "FAKTA KUNCI", 22, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(9.2), Inches(2.6), Inches(1.5), Inches(0.05), WHITE)

facts = [
    ("10+", "Tahun Pengalaman"),
    ("150+", "Proyek Selesai"),
    ("100+", "Klien Puas"),
    ("dormakaba", "Partner Resmi"),
    ("100%", "Kepuasan Klien")
]

y_pos = 3
for num, label in facts:
    add_textbox(slide, Inches(9.2), Inches(y_pos), Inches(3), Inches(0.5), 
                num, 28, True, WHITE, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(9.2), Inches(y_pos + 0.4), Inches(3), Inches(0.4), 
                label, 16, False, RGBColor(0xFF, 0xD4, 0xD4), PP_ALIGN.LEFT)
    y_pos += 0.85

# ============================================
# SLIDE 3: VISION & MISSION
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), 
            "VISI & MISI PERUSAHAAN", 32, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1), Inches(0.08), WHITE)

# Vision Box
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), fill_color=RED_LIGHT)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(0.12), Inches(2.3), DARK_RED)

add_textbox(slide, Inches(1.3), Inches(2), Inches(4.8), Inches(0.5), 
            "VISI KAMI", 24, True, DARK_RED, PP_ALIGN.LEFT)

add_textbox(slide, Inches(1.3), Inches(2.6), Inches(4.8), Inches(1.3), 
            "Menjadi perusahaan general kontraktor terdepan yang dikenal dengan kualitas, integritas, dan inovasi dalam industri konstruksi Indonesia.", 
            18, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Mission Box
add_shape(slide, Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.3), fill_color=RED_LIGHT)
add_accent_bar(slide, Inches(6.9), Inches(1.8), Inches(0.12), Inches(2.3), DARK_RED)

add_textbox(slide, Inches(7.4), Inches(2), Inches(4.8), Inches(0.5), 
            "MISI KAMI", 24, True, DARK_RED, PP_ALIGN.LEFT)

missions_text = ("• Memberikan layanan konstruksi berkualitas tinggi\n"
                 "• Menyediakan solusi pintu otomatis dormakaba\n"
                 "• Membangun hubungan jangka panjang\n"
                 "• Berkontribusi pada infrastruktur Indonesia")
add_textbox(slide, Inches(7.4), Inches(2.5), Inches(4.8), Inches(1.4), 
            missions_text, 16, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Core Values Section
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5), 
            "NILAI-NILAI INTI", 24, True, DARK_GRAY, PP_ALIGN.LEFT)

values = [
    ("Kualitas", "Standar internasional dalam setiap proyek"),
    ("Integritas", "Kejujuran dan transparansi dalam bekerja"),
    ("Inovasi", "Solusi modern untuk kebutuhan konstruksi"),
    ("Kepuasan Klien", "Fokus utama dalam setiap pelayanan")
]

x_pos = 0.8
for i, (title, desc) in enumerate(values):
    add_shape(slide, Inches(x_pos), Inches(5.1), Inches(2.8), Inches(1.8), fill_color=VERY_LIGHT_GRAY)
    add_accent_bar(slide, Inches(x_pos), Inches(5.1), Inches(2.8), Inches(0.08), DARK_RED)
    
    add_textbox(slide, Inches(x_pos + 0.2), Inches(5.3), Inches(2.4), Inches(0.4), 
                title, 20, True, DARK_RED, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x_pos + 0.2), Inches(5.8), Inches(2.4), Inches(0.8), 
                desc, 14, False, MEDIUM_GRAY, PP_ALIGN.CENTER)
    
    x_pos += 3.1

# ============================================
# SLIDE 4: OUR SERVICES OVERVIEW
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), 
            "LAYANAN UTAMA KAMI", 32, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1), Inches(0.08), WHITE)

# Service boxes
services = [
    ("General Construction", "Pembangunan baru, renovasi, dan rehabilitasi dengan standar kualitas tertinggi dan tim profesional berpengalaman"),
    ("dormakaba Automatic Doors", "Instalasi, maintenance, dan perbaikan sistem pintu otomatis berkualitas tinggi untuk gedung komersial dan residensial"),
    ("Interior & Renovation", "Jasa renovasi dan desain interior profesional untuk rumah, kantor, hotel, dan ruang komersial")
]

for i, (title, desc) in enumerate(services):
    y = 1.8 + (i * 1.7)
    
    # Service box
    add_shape(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(1.4), fill_color=VERY_LIGHT_GRAY)
    add_accent_bar(slide, Inches(0.8), Inches(y), Inches(0.12), Inches(1.4), DARK_RED)
    
    # Number circle
    add_circle(slide, Inches(1.3), Inches(y + 0.35), Inches(0.7), DARK_RED)
    add_textbox(slide, Inches(1.3), Inches(y + 0.38), Inches(0.7), Inches(0.6), 
                str(i + 1), 24, True, WHITE, PP_ALIGN.CENTER)
    
    # Title and description
    add_textbox(slide, Inches(2.3), Inches(y + 0.15), Inches(9.5), Inches(0.4), 
                title, 22, True, DARK_GRAY, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(2.3), Inches(y + 0.6), Inches(9.5), Inches(0.7), 
                desc, 16, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# ============================================
# SLIDE 5: GENERAL CONSTRUCTION DETAIL
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), 
            "GENERAL CONSTRUCTION", 32, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1), Inches(0.08), WHITE)

# Left column
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5), fill_color=VERY_LIGHT_GRAY)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(0.12), Inches(5), DARK_RED)

add_textbox(slide, Inches(1.3), Inches(2), Inches(5), Inches(0.5), 
            "Layanan Konstruksi Menyeluruh", 24, True, DARK_RED, PP_ALIGN.LEFT)

add_textbox(slide, Inches(1.3), Inches(2.7), Inches(5), Inches(3.8), 
            "Kami menyediakan layanan konstruksi komprehensif untuk berbagai kebutuhan:\n\n"
            "✓ Pembangunan Baru\n"
            "  - Gedung komersial dan perkantoran\n"
            "  - Fasilitas industri dan manufaktur\n"
            "  - Hotel dan hospitality\n\n"
            "✓ Renovasi Besar\n"
            "  - Renovasi total bangunan\n"
            "  - Upgrade sistem dan infrastruktur\n"
            "  - Modernisasi fasilitas\n\n"
            "✓ Manajemen Proyek\n"
            "  - Perencanaan dan pengawasan\n"
            "  - Pengendalian kualitas dan biaya\n"
            "  - Penyelesaian tepat waktu", 
            18, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Right column - Why choose us
add_shape(slide, Inches(7.1), Inches(1.8), Inches(5.4), Inches(5), fill_color=DARK_RED)

add_textbox(slide, Inches(7.5), Inches(2.1), Inches(4.5), Inches(0.5), 
            "MENGAPA MEMILIH KAMI", 22, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(7.5), Inches(2.6), Inches(2), Inches(0.05), WHITE)

reasons = [
    "Tim profesional berpengalaman lebih dari 10 tahun",
    "Standar kualitas internasional",
    "Material berkualitas tinggi",
    "Pengawasan ketat di setiap tahap",
    "Tepat waktu dan sesuai budget",
    "Garansi hasil pekerjaan",
    "Layanan purna jual responsif",
    "Harga kompetitif dan transparan"
]

y_pos = 2.9
for reason in reasons:
    add_textbox(slide, Inches(7.5), Inches(y_pos), Inches(4.5), Inches(0.5), 
                "✓  " + reason, 16, False, WHITE, PP_ALIGN.LEFT)
    y_pos += 0.55

# ============================================
# SLIDE 6: AUTOMATIC DOOR SOLUTIONS
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), 
            "dormakaba AUTOMATIC DOOR SOLUTIONS", 32, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1), Inches(0.08), WHITE)

# Left - dormakaba partnership info
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5), fill_color=VERY_LIGHT_GRAY)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(0.12), Inches(5), DARK_RED)

add_textbox(slide, Inches(1.3), Inches(2), Inches(5), Inches(0.5), 
            "Partner Resmi dormakaba", 24, True, DARK_RED, PP_ALIGN.LEFT)

add_textbox(slide, Inches(1.3), Inches(2.7), Inches(5), Inches(3.8), 
            "Sebagai partner resmi dormakaba, kami menyediakan:\n\n"
            "✓ Instalasi Pintu Otomatis\n"
            "  - Sliding doors\n"
            "  - Swing doors\n"
            "  - Revolving doors\n"
            "  - Security access systems\n\n"
            "✓ Maintenance & Perawatan\n"
            "  - Perawatan berkala\n"
            "  - Inspections rutin\n"
            "  - Preventive maintenance\n\n"
            "✓ Perbaikan & Support\n"
            "  - Emergency repair 24/7\n"
            "  - Spare parts original\n"
            "  - Technical support", 
            18, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Right - Benefits
add_shape(slide, Inches(7.1), Inches(1.8), Inches(5.4), Inches(5), fill_color=DARK_RED)

add_textbox(slide, Inches(7.5), Inches(2.1), Inches(4.5), Inches(0.5), 
            "KEUNGGULAN dormakaba", 22, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(7.5), Inches(2.6), Inches(2), Inches(0.05), WHITE)

benefits = [
    "Brand global terpercaya sejak 1868",
    "Teknologi terdepan di industri",
    "Hemat energi hingga 30%",
    "Aksesibilitas untuk semua orang",
    "Keamanan tingkat tinggi",
    "Desain modern dan elegan",
    "Daya tahan dan keandalan",
    "Garansi resmi dormakaba",
    "Layanan after-sales premium"
]

y_pos = 2.9
for benefit in benefits:
    add_textbox(slide, Inches(7.5), Inches(y_pos), Inches(4.5), Inches(0.5), 
                "✓  " + benefit, 16, False, WHITE, PP_ALIGN.LEFT)
    y_pos += 0.48

# ============================================
# SLIDE 7: CLIENTS / TRUSTED PARTNERS
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), 
            "DIPERCAYA OLEH PERUSAHAAN TERKEMUKA", 32, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1), Inches(0.08), WHITE)

# Intro text
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6), 
            "Kami telah dipercaya oleh lebih dari 100 perusahaan di berbagai industri", 20, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Clients list in grid
clients = [
    "DoubleTree by Hilton",
    "Horison Hotels",
    "ITC Group",
    "LG Innotek",
    "Sasa Foods",
    "Abipraya",
    "Astra International",
    "BSJ Group"
]

x_pos = 0.8
y_pos = 2.5
for i, client in enumerate(clients):
    if i == 4:
        x_pos = 0.8
        y_pos = 3.9
    
    add_shape(slide, Inches(x_pos), Inches(y_pos), Inches(2.8), Inches(1), fill_color=VERY_LIGHT_GRAY)
    add_accent_bar(slide, Inches(x_pos), Inches(y_pos), Inches(2.8), Inches(0.06), DARK_RED)
    add_textbox(slide, Inches(x_pos + 0.1), Inches(y_pos + 0.2), Inches(2.6), Inches(0.5), 
                client, 18, True, DARK_GRAY, PP_ALIGN.CENTER)
    
    x_pos += 3.1

# Bottom section - Industries served
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5), 
            "INDUSTRI YANG KAMI LAYANI", 24, True, DARK_GRAY, PP_ALIGN.LEFT)

industries = ["Hospitality & Hotel", "Manufaktur & Industri", "Retail & Shopping Mall", "Perkantoran & Komersial", "Infrastruktur & Pemerintah"]
x_pos = 0.8
for industry in industries:
    add_shape(slide, Inches(x_pos), Inches(5.8), Inches(2.3), Inches(0.7), fill_color=RED_LIGHT)
    add_textbox(slide, Inches(x_pos + 0.1), Inches(5.9), Inches(2.1), Inches(0.5), 
                industry, 16, True, DARK_RED, PP_ALIGN.CENTER)
    x_pos += 2.5

# ============================================
# SLIDE 8: WHY CHOOSE US
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), 
            "MENGAPA MEMILIH PT CAHAYA DIMENSI BUMI?", 32, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1), Inches(0.08), WHITE)

# Reasons in 2 columns
reasons_left = [
    ("Pengalaman 10+ Tahun", "Lebih dari satu dekara pengalaman dalam industri konstruksi Indonesia"),
    ("Partner Resmi dormakaba", "Akses langsung ke produk dan dukungan teknis dari dormakaba"),
    ("Tim Profesional", "Tenaga ahli bersertifikat dengan kompetensi teruji"),
    ("Kualitas Terjamin", "Standar internasional dengan pengawasan ketat")
]

reasons_right = [
    ("Harga Kompetitif", "Penawaran terbaik dengan kualitas premium"),
    ("Tepat Waktu", "Komitmen penyelesaian sesuai jadwal"),
    ("Garansi Pekerjaan", "Jaminan hasil pekerjaan dan after-sales support"),
    ("100+ Klien Puas", "Track record terpercaya di berbagai industri")
]

y_pos = 1.8
for i, (title, desc) in enumerate(reasons_left):
    add_shape(slide, Inches(0.8), Inches(y_pos), Inches(5.6), Inches(1.2), fill_color=VERY_LIGHT_GRAY)
    add_accent_bar(slide, Inches(0.8), Inches(y_pos), Inches(0.1), Inches(1.2), DARK_RED)
    
    add_textbox(slide, Inches(1.3), Inches(y_pos + 0.1), Inches(4.8), Inches(0.35), 
                "✓  " + title, 18, True, DARK_RED, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1.3), Inches(y_pos + 0.5), Inches(4.8), Inches(0.6), 
                desc, 14, False, MEDIUM_GRAY, PP_ALIGN.LEFT)
    
    y_pos += 1.4

y_pos = 1.8
for i, (title, desc) in enumerate(reasons_right):
    add_shape(slide, Inches(6.9), Inches(y_pos), Inches(5.6), Inches(1.2), fill_color=VERY_LIGHT_GRAY)
    add_accent_bar(slide, Inches(6.9), Inches(y_pos), Inches(0.1), Inches(1.2), DARK_RED)
    
    add_textbox(slide, Inches(7.4), Inches(y_pos + 0.1), Inches(4.8), Inches(0.35), 
                "✓  " + title, 18, True, DARK_RED, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(7.4), Inches(y_pos + 0.5), Inches(4.8), Inches(0.6), 
                desc, 14, False, MEDIUM_GRAY, PP_ALIGN.LEFT)
    
    y_pos += 1.4

# Bottom CTA
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.5), 
            "Konsultasi GRATIS untuk kebutuhan proyek Anda — Hubungi kami sekarang!", 20, True, WHITE, PP_ALIGN.CENTER)

# ============================================
# SLIDE 9: CONTACT INFORMATION
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Header
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.8), 
            "HUBUNGI KAMI", 32, True, WHITE, PP_ALIGN.LEFT)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(1), Inches(0.08), WHITE)

# Contact details
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.6), 
            "Siap membantu mewujudkan proyek Anda. Hubungi kami untuk konsultasi gratis.", 20, False, MEDIUM_GRAY, PP_ALIGN.LEFT)

# Contact boxes
contact_info = [
    ("Kantor", "Jakarta, Indonesia"),
    ("Telepon/WhatsApp", "+62 851-7171-1375"),
    ("Email", "info@cahayadimensibumi.com"),
    ("Jam Operasional", "Senin - Jumat: 09:00 - 18:00\nSabtu: 09:00 - 14:00")
]

for i, (label, value) in enumerate(contact_info):
    x = 0.8 + (i * 3.1)
    
    add_shape(slide, Inches(x), Inches(2.8), Inches(2.8), Inches(2.5), fill_color=VERY_LIGHT_GRAY)
    add_accent_bar(slide, Inches(x), Inches(2.8), Inches(2.8), Inches(0.1), DARK_RED)
    
    add_textbox(slide, Inches(x + 0.2), Inches(3.1), Inches(2.4), Inches(0.4), 
                label, 18, True, DARK_RED, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x + 0.2), Inches(3.7), Inches(2.4), Inches(1.2), 
                value, 18, False, DARK_GRAY, PP_ALIGN.CENTER)

# WhatsApp CTA
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), fill_color=DARK_RED)
add_textbox(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.5), 
            "WhatsApp Kami Sekarang", 28, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5), 
            "Konsultasi GRATIS | Respon Cepat | Penawaran Kompetitif", 18, False, RGBColor(0xFF, 0xD4, 0xD4), PP_ALIGN.CENTER)

# ============================================
# SLIDE 10: CLOSING / THANK YOU
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Full screen red background
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill_color=DARK_RED)

# Decorative elements
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0xEF, 0x44, 0x44))
add_circle(slide, Inches(11), Inches(5.5), Inches(3.5), RGBColor(0xEF, 0x44, 0x44))

# Thank you text
add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.2), 
            "TERIMA KASIH", 56, True, WHITE, PP_ALIGN.CENTER)

add_shape(slide, Inches(5), Inches(3.4), Inches(3.333), Inches(0.08), fill_color=WHITE)

add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8), 
            "Atas kepercayaan dan perhatian Bapak/Ibu", 28, False, RGBColor(0xFF, 0xD4, 0xD4), PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8), 
            "PT CAHAYA DIMENSI BUMI", 32, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.6), 
            "General Contractor & Automatic Door Solutions", 20, False, RGBColor(0xFF, 0xD4, 0xD4), PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.6), 
            "WhatsApp: +62 851-7171-1375 | info@cahayadimensibumi.com", 18, False, WHITE, PP_ALIGN.CENTER)

# Save
output_path = r'C:\xampp\htdocs\web\Cahaya_Dimensi_Bumi_Company_Profile.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")